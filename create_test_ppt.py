from pptx import Presentation
from pptx.util import Inches

def create_test_ppt(filename):
    prs = Presentation()

    # Slide 1
    slide1_layout = prs.slide_layouts[5]
    slide1 = prs.slides.add_slide(slide1_layout)
    title1 = slide1.shapes.title
    title1.text = "This is a test presentation."

    # Slide 2
    slide2_layout = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(slide2_layout)
    title2 = slide2.shapes.title
    title2.text = "Topic 1: Python"
    content2 = slide2.placeholders[1]
    content2.text = "Python is a popular programming language."

    # Slide 3
    slide3_layout = prs.slide_layouts[1]
    slide3 = prs.slides.add_slide(slide3_layout)
    title3 = slide3.shapes.title
    title3.text = "Topic 2: Streamlit"
    content3 = slide3.placeholders[1]
    content3.text = "Streamlit is a great tool for building web apps."

    prs.save(filename)

if __name__ == "__main__":
    create_test_ppt("test.pptx")
