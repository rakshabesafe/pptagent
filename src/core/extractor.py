import os
from pptx import Presentation
from PIL import Image
import pytesseract

def extract_text_from_slide(slide):
    text = ""
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            text += shape.text + "\n"
    return text

def extract_images_from_slide(slide, output_dir):
    image_paths = []
    for i, shape in enumerate(slide.shapes):
        if shape.shape_type == 13:  # Picture
            image = shape.image
            image_bytes = image.blob
            image_filename = f"slide_{slide.slide_id}_image_{i}.{image.ext}"
            image_path = os.path.join(output_dir, image_filename)
            with open(image_path, "wb") as f:
                f.write(image_bytes)
            image_paths.append(image_path)
    return image_paths

def ocr_image(image_path):
    try:
        return pytesseract.image_to_string(Image.open(image_path))
    except Exception as e:
        return f"Error during OCR: {e}"

def extract_metadata(file_path):
    presentation = Presentation(file_path)
    metadata = []
    image_dir = "images"
    if not os.path.exists(image_dir):
        os.makedirs(image_dir)

    for i, slide in enumerate(presentation.slides):
        text = extract_text_from_slide(slide)
        image_paths = extract_images_from_slide(slide, image_dir)
        ocr_texts = [ocr_image(p) for p in image_paths]

        metadata.append({
            "slide_number": i + 1,
            "text": text,
            "image_paths": image_paths,
            "ocr_text": "\n".join(ocr_texts)
        })

    return metadata
